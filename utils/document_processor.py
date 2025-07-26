import logging
import io
from typing import Union, Optional
import PyPDF2
import pdfplumber
from docx import Document
import openpyxl
from bs4 import BeautifulSoup


class DocumentProcessor:
    """Utility class for processing various document types"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def process_document(self, content: bytes, file_type: str) -> str:
        """
        Process document content based on file type
        
        Args:
            content: Raw document content as bytes
            file_type: File extension (pdf, doc, docx, etc.)
            
        Returns:
            Extracted text content as string
        """
        try:
            if file_type.lower() == 'pdf':
                return self.process_pdf(content)
            elif file_type.lower() in ['doc', 'docx']:
                return self.process_word_document(content)
            elif file_type.lower() in ['xls', 'xlsx']:
                return self.process_excel_document(content)
            elif file_type.lower() in ['ppt', 'pptx']:
                return self.process_powerpoint_document(content)
            elif file_type.lower() == 'txt':
                return self.process_text_document(content)
            elif file_type.lower() == 'rtf':
                return self.process_rtf_document(content)
            elif file_type.lower() in ['html', 'htm']:
                return self.process_html_document(content)
            else:
                self.logger.warning(f"Unsupported file type: {file_type}")
                return ""
                
        except Exception as e:
            self.logger.error(f"Error processing {file_type} document: {str(e)}")
            return ""
    
    def process_pdf(self, content: bytes) -> str:
        """Process PDF documents"""
        text_content = ""
        
        try:
            # Try pdfplumber first (better text extraction)
            with io.BytesIO(content) as pdf_buffer:
                with pdfplumber.open(pdf_buffer) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_content += page_text + "\n"
        except Exception as e:
            self.logger.warning(f"pdfplumber failed, trying PyPDF2: {str(e)}")
            
            # Fallback to PyPDF2
            try:
                with io.BytesIO(content) as pdf_buffer:
                    pdf_reader = PyPDF2.PdfReader(pdf_buffer)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_content += page_text + "\n"
            except Exception as e2:
                self.logger.error(f"Both PDF processors failed: {str(e2)}")
        
        return text_content.strip()
    
    def process_word_document(self, content: bytes) -> str:
        """Process Word documents (DOC/DOCX)"""
        try:
            with io.BytesIO(content) as doc_buffer:
                doc = Document(doc_buffer)
                text_content = []
                
                for paragraph in doc.paragraphs:
                    text_content.append(paragraph.text)
                
                # Also extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text_content.append(cell.text)
                
                return "\n".join(text_content)
                
        except Exception as e:
            self.logger.error(f"Error processing Word document: {str(e)}")
            return ""
    
    def process_excel_document(self, content: bytes) -> str:
        """Process Excel documents (XLS/XLSX)"""
        try:
            with io.BytesIO(content) as excel_buffer:
                workbook = openpyxl.load_workbook(excel_buffer, data_only=True)
                text_content = []
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_content.append(f"Sheet: {sheet_name}")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text:
                            text_content.append("\t".join(row_text))
                
                return "\n".join(text_content)
                
        except Exception as e:
            self.logger.error(f"Error processing Excel document: {str(e)}")
            return ""
    
    def process_powerpoint_document(self, content: bytes) -> str:
        """Process PowerPoint documents (PPT/PPTX)"""
        try:
            from pptx import Presentation
            
            with io.BytesIO(content) as ppt_buffer:
                prs = Presentation(ppt_buffer)
                text_content = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_content.append(f"Slide {slide_num}:")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content.append(shape.text)
                    
                    text_content.append("")  # Empty line between slides
                
                return "\n".join(text_content)
            
        except ImportError:
            self.logger.warning("python-pptx not installed. Cannot process PowerPoint files.")
            return "PowerPoint document detected - python-pptx required for text extraction"
        except Exception as e:
            self.logger.error(f"Error processing PowerPoint document: {str(e)}")
            return ""
    
    def process_text_document(self, content: bytes) -> str:
        """Process plain text documents"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use errors='ignore'
            return content.decode('utf-8', errors='ignore')
            
        except Exception as e:
            self.logger.error(f"Error processing text document: {str(e)}")
            return ""
    
    def process_rtf_document(self, content: bytes) -> str:
        """Process RTF documents"""
        try:
            # Basic RTF processing - strip RTF commands
            text = content.decode('utf-8', errors='ignore')
            
            # Remove RTF control sequences (very basic)
            import re
            # Remove RTF control words
            text = re.sub(r'\\[a-z]+\d*', '', text)
            # Remove curly braces
            text = re.sub(r'[{}]', '', text)
            # Clean up multiple spaces and newlines
            text = re.sub(r'\s+', ' ', text)
            
            return text.strip()
            
        except Exception as e:
            self.logger.error(f"Error processing RTF document: {str(e)}")
            return ""
    
    def process_html_document(self, content: bytes) -> str:
        """Process HTML documents"""
        try:
            html_content = content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception as e:
            self.logger.error(f"Error processing HTML document: {str(e)}")
            return ""
    
    def extract_metadata(self, content: bytes, file_type: str) -> dict:
        """Extract metadata from documents"""
        metadata = {}
        
        try:
            if file_type.lower() == 'pdf':
                with io.BytesIO(content) as pdf_buffer:
                    pdf_reader = PyPDF2.PdfReader(pdf_buffer)
                    if pdf_reader.metadata:
                        metadata.update({
                            'title': pdf_reader.metadata.get('/Title', ''),
                            'author': pdf_reader.metadata.get('/Author', ''),
                            'subject': pdf_reader.metadata.get('/Subject', ''),
                            'creator': pdf_reader.metadata.get('/Creator', ''),
                            'producer': pdf_reader.metadata.get('/Producer', ''),
                            'creation_date': pdf_reader.metadata.get('/CreationDate', ''),
                            'modification_date': pdf_reader.metadata.get('/ModDate', ''),
                        })
                    metadata['page_count'] = len(pdf_reader.pages)
                    
            elif file_type.lower() in ['doc', 'docx']:
                with io.BytesIO(content) as doc_buffer:
                    doc = Document(doc_buffer)
                    if doc.core_properties:
                        metadata.update({
                            'title': doc.core_properties.title or '',
                            'author': doc.core_properties.author or '',
                            'subject': doc.core_properties.subject or '',
                            'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                            'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
                        })
                        
        except Exception as e:
            self.logger.error(f"Error extracting metadata from {file_type}: {str(e)}")
            
        return metadata